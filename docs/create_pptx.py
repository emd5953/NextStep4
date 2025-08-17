from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import sys

def create_poster():
    try:
        prs = Presentation()
        
        # Set slide dimensions (48x36 inches)
        prs.slide_width = Inches(48)
        prs.slide_height = Inches(36)
        
        # Add blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add header shape
        header = slide.shapes.add_shape(1, 0, 0, Inches(48), Inches(4))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(45, 20, 65)  # #2D1441
        header.line.color.rgb = RGBColor(26, 11, 38)  # #1a0b26
        
        # Add title
        title = header.text_frame
        title.text = "NEXTSTEP: Swipe-Based Job Matching App"
        title.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.paragraphs[0].font.size = Pt(60)
        title.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.paragraphs[0].font.bold = True
        
        # Add subtitle
        subtitle = slide.shapes.add_textbox(0, Inches(1.5), Inches(48), Inches(1))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = "Your next career move, simplified."
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(204, 204, 204)
        
        # Add authors
        authors = slide.shapes.add_textbox(0, Inches(2.5), Inches(48), Inches(1))
        authors_frame = authors.text_frame
        authors_frame.text = "Andrew Nguyen • Enrin Debbarma • Soleyana Abera • Bryan R Mathews"
        authors_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        authors_frame.paragraphs[0].font.size = Pt(30)
        authors_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add affiliation
        affiliation = slide.shapes.add_textbox(0, Inches(3.2), Inches(48), Inches(1))
        affiliation_frame = affiliation.text_frame
        affiliation_frame.text = "Engineering: Computer Science Students"
        affiliation_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        affiliation_frame.paragraphs[0].font.size = Pt(24)
        affiliation_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Define section dimensions
        section_width = Inches(15)
        section_height = Inches(28)
        start_y = Inches(5)
        
        # Create sections with titles
        sections = [
            ("Purpose", [
                "[INSERT ABOUT_NEXTSTEP.PNG IMAGE HERE]",
                "",
                "NextStep is designed to make job searching easier and more interactive. With a swipe-based interface, users can browse jobs effortlessly and apply with just one click.",
                "",
                "Key Features:",
                "",
                "• Swipe-based job browsing for an engaging experience",
                "• AI-powered personalized job recommendations",
                "• One-Click Apply for faster job applications",
                "• Application tracking to manage job searches efficiently"
            ]),
            ("System Architecture", [
                "Technology Stack:",
                "",
                "• Frontend: React.js & React Native",
                "• Backend: Node.js with Express.js",
                "• Database: Cloud-based MongoDB",
                "• Authentication: Google OAuth",
                "• Real-time Updates: Firebase",
                "",
                "[INSERT SIGNUP_SCREEN.PNG IMAGE HERE]",
                "",
                "User Authentication:",
                "",
                "• Secure phone verification via Twilio",
                "• Email validation system",
                "• Password encryption",
                "• Role-based access control"
            ]),
            ("Security & Optimization", [
                "Data Privacy & Security:",
                "",
                "We take privacy seriously. Your data is securely stored and encrypted to ensure a safe job-hunting experience.",
                "",
                "• End-to-end encryption",
                "• GDPR compliance",
                "• Regular security audits",
                "• Secure cloud storage",
                "",
                "Performance Optimization:",
                "",
                "NextStep helps users track and analyze job applications with built-in insights, so they can refine their job search strategy.",
                "",
                "• Real-time analytics",
                "• Personalized insights",
                "• Application success tracking",
                "• Job market trends analysis"
            ])
        ]
        
        # Add sections
        for i, (title, content) in enumerate(sections):
            x = Inches(1 + i * 16)
            
            # Add section background
            section = slide.shapes.add_shape(1, x, start_y, section_width, section_height)
            section.fill.solid()
            section.fill.fore_color.rgb = RGBColor(255, 255, 255)
            section.line.color.rgb = RGBColor(217, 217, 217)
            
            # Add section title
            title_shape = slide.shapes.add_shape(1, x, start_y, section_width, Inches(1))
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(45, 20, 65)  # #2D1441
            title_shape.line.color.rgb = RGBColor(26, 11, 38)  # #1a0b26
            title_frame = title_shape.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.size = Pt(36)
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            title_frame.paragraphs[0].font.bold = True
            
            # Add content
            content_box = slide.shapes.add_textbox(x + Inches(0.5), start_y + Inches(1.5), 
                                                 section_width - Inches(1), section_height - Inches(2))
            add_content_to_shape(content_box, content)
            
            # Try to add images
            if title == "Purpose":
                try:
                    img_path = os.path.join(os.getcwd(), "about_nextstep.png")
                    if os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, x + Inches(1), start_y + Inches(2), 
                                              width=section_width - Inches(2))
                except Exception as e:
                    print(f"Could not add about_nextstep.png: {str(e)}")
            
            elif title == "System Architecture":
                try:
                    img_path = os.path.join(os.getcwd(), "signup_screen.png")
                    if os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, x + Inches(1), start_y + Inches(10), 
                                              width=section_width - Inches(2))
                except Exception as e:
                    print(f"Could not add signup_screen.png: {str(e)}")
        
        # Add footer
        footer = slide.shapes.add_shape(1, 0, Inches(33), Inches(48), Inches(3))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(45, 20, 65)  # #2D1441
        footer.line.color.rgb = RGBColor(26, 11, 38)  # #1a0b26
        
        footer_text = footer.text_frame
        footer_text.text = "Contact: info@nextstep4.com | www.nextstep4.com | (123) 456-7890"
        footer_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        footer_text.paragraphs[0].font.size = Pt(24)
        footer_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Try different save locations
        save_locations = [
            os.path.join(os.path.expanduser("~"), "Desktop", "NextStep_Poster.pptx"),
            os.path.join(os.path.expanduser("~"), "Documents", "NextStep_Poster.pptx"),
            os.path.join(os.getcwd(), "NextStep_Poster.pptx")
        ]
        
        saved = False
        for save_path in save_locations:
            try:
                prs.save(save_path)
                print(f"Presentation saved successfully to: {save_path}")
                saved = True
                break
            except PermissionError:
                print(f"Could not save to {save_path}, trying next location...")
                continue
        
        if not saved:
            print("Error: Could not save the presentation to any location.")
            print("Please try running the script with administrator privileges or choose a different save location.")
            sys.exit(1)
            
    except Exception as e:
        print(f"An error occurred: {str(e)}")
        sys.exit(1)

def add_content_to_shape(shape, content):
    text_frame = shape.text_frame
    first = True
    
    for line in content:
        if not first:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
            first = False
        
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(51, 51, 51)
        
        if line.startswith('•'):
            p.font.color.rgb = RGBColor(45, 20, 65)  # Purple for bullet points
        elif line.startswith('[INSERT'):
            p.font.color.rgb = RGBColor(128, 128, 128)
            p.font.italic = True
        elif line.endswith(':'):
            p.font.bold = True
            p.font.size = Pt(28)

if __name__ == "__main__":
    create_poster() 