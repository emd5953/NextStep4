from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_poster():
    prs = Presentation()
    
    # Set slide size to 48" x 36"
    prs.slide_width = Inches(48)
    prs.slide_height = Inches(36)
    
    # Add slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    
    # Add header
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(48), Inches(6))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(45, 20, 65)  # #2D1441
    
    # Add title
    title = header.text_frame
    title.text = "NEXTSTEP: Swipe-Based Job Matching App"
    title.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.paragraphs[0].font.size = Pt(72)
    title.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add subtitle
    subtitle = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(48), Inches(1))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "PENN STATE CAPSTONE 2025"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(48)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add team names
    team = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(48), Inches(1))
    team_frame = team.text_frame
    team_frame.text = "Andrew Nguyen • Enrin Debbarma • Soleyana Abera • Bryan R Mathews"
    team_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    team_frame.paragraphs[0].font.size = Pt(32)
    team_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Create three columns
    column_width = Inches(15)
    column_height = Inches(26)
    spacing = Inches(1.5)
    start_y = Inches(7)
    
    # Left Column
    left_col = slide.shapes.add_shape(1, Inches(0.5), start_y, column_width, column_height)
    left_col.fill.solid()
    left_col.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_col.line.color.rgb = RGBColor(229, 231, 235)
    
    # Middle Column
    mid_col = slide.shapes.add_shape(1, Inches(16.5), start_y, column_width, column_height)
    mid_col.fill.solid()
    mid_col.fill.fore_color.rgb = RGBColor(255, 255, 255)
    mid_col.line.color.rgb = RGBColor(229, 231, 235)
    
    # Right Column
    right_col = slide.shapes.add_shape(1, Inches(32.5), start_y, column_width, column_height)
    right_col.fill.solid()
    right_col.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_col.line.color.rgb = RGBColor(229, 231, 235)
    
    # Add content to columns
    add_content_to_shape(left_col, "Abstract", abstract_content())
    add_content_to_shape(mid_col, "System Architecture", system_arch_content())
    add_content_to_shape(right_col, "Security & Performance", security_content())
    
    # Add footer
    footer = slide.shapes.add_shape(1, Inches(0), Inches(34), Inches(48), Inches(2))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(45, 20, 65)  # #2D1441
    
    footer_text = footer.text_frame
    footer_text.text = "Contact: info@nextstep4.com | www.nextstep4.com | (123) 456-7890"
    footer_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_text.paragraphs[0].font.size = Pt(32)
    footer_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Save presentation
    prs.save('NextStep_Poster_Final.pptx')

def add_content_to_shape(shape, title, content):
    text_frame = shape.text_frame
    text_frame.text = title
    text_frame.paragraphs[0].font.size = Pt(36)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(45, 20, 65)
    
    for line in content:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(75, 85, 99)

def abstract_content():
    return [
        "\nNextStep revolutionizes the job search experience by combining the intuitive simplicity of swipe-based interfaces with powerful AI-driven matching technology. Our platform addresses the fundamental challenges in today's job market by creating seamless connections between job seekers and employers.",
        "\nOur AI-powered platform learns from each interaction, creating a dynamic matching system that gets smarter with every swipe. By analyzing user behavior, skills, and preferences, NextStep delivers precisely targeted job opportunities to candidates while providing employers with pre-qualified, interested applicants.",
        "\nKey Innovation Points:",
        "• AI-Powered Job Matching",
        "• Intuitive Swipe Interface",
        "• Real-time Application Tracking",
        "• Secure Data Management"
    ]

def system_arch_content():
    return [
        "\nTechnology Stack:",
        "\nFrontend:",
        "• React.js (Web Application)",
        "• React Native (Mobile Apps)",
        "\nBackend:",
        "• Node.js with Express.js",
        "• Cloud-based MongoDB",
        "\nExternal Services:",
        "• Google OAuth Authentication",
        "• Twilio Phone Verification",
        "• Firebase Real-time Updates"
    ]

def security_content():
    return [
        "\nData Protection:",
        "• End-to-end encryption",
        "• GDPR compliance",
        "• Regular security audits",
        "• Secure data storage",
        "\nSystem Performance:",
        "• 100,000+ concurrent users",
        "• 95% pages load < 2 seconds",
        "• 99.9% uptime guarantee",
        "• Mobile-optimized experience"
    ]

if __name__ == "__main__":
    create_poster() 